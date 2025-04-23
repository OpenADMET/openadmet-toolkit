import logging
from pathlib import Path
from tempfile import TemporaryDirectory

import cairosvg
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from rdkit import Chem
from rdkit.Chem import Draw

logger = logging.getLogger(__name__)


try:
    from IPython.core.display import SVG  # noqa: F401
    from rdkit.Chem.Draw import IPythonConsole, rdMolDraw2D  # noqa: F401
except ImportError:
    IPythonConsole = None


def rdkit_draw_in_ipython_mode() -> bool:
    """
    Check if RDKit is in IPython mode.

    Returns
    -------
    bool
        True if RDKit is in IPython mode, False otherwise.

    """
    return IPythonConsole is not None


def create_presentation(image_paths, output_file) -> None:
    # Create a presentation object
    prs = Presentation()

    for image_path in image_paths:
        # Add a slide with a blank layout
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 6 = blank layout

        # Add the image to the slide
        left = Inches(1)
        top = Inches(0)
        width = Inches(8)
        slide.shapes.add_picture(image_path, left, top, width=width)

    # Save the presentation
    prs.save(output_file)
    logger.info(f"Presentation saved to {output_file}")


def make_pptx_from_molecule_data(
    dataframe: pd.DataFrame,
    output_file: str | Path,
    smiles_col="SMILES",
    mols_per_slide: int = 25,
    keep_images=False,
    image_dir: str | Path = "images",
    legend_columns=None,
) -> None:
    """
    This function takes a dataframe with molecule data and creates a pptx file with the molecules as images.

    Parameters
    ----------
    dataframe : pandas.DataFrame
        A dataframe with molecule data.
    output_file : str
        The name of the output pptx file.
    smiles_col : str, Optional, default: "SMILES"
        The name of the column in the dataframe that contains the SMILES strings of the molecules.
    keep_images : bool, Optional, default: False
        If True, the images of the molecules will be saved as .png files.
    image_dir : str, Optional, default: "images"
        The directory where the images will be saved.
    legend_columns : list, Optional, default: None
        A list of column names from the dataframe that will be used as legend for the images.

    Returns
    -------
    None

    """
    total_mols = len(dataframe)
    num_slides = np.ceil(total_mols / mols_per_slide).astype(int)
    chunks = np.array_split(dataframe, num_slides)

    logger.debug(f"Creating {num_slides} slides with {mols_per_slide} molecules each")

    if keep_images:
        Path(image_dir).mkdir(exist_ok=True)
    else:
        # setup temp_dir
        temp_dir = TemporaryDirectory()
        image_dir = temp_dir.name

    img_path = Path(image_dir)

    # take square root of mols_per_slide to get the number of mols per row
    # by default will make a 5x5 grid

    mols_per_row = int(np.sqrt(mols_per_slide).astype(int))

    svgnames = []
    pngnames = []

    if legend_columns:
        for col in legend_columns:
            if col not in dataframe.columns:
                raise ValueError(f"Column {col} not found in dataframe")

    if smiles_col not in dataframe.columns:
        raise ValueError(f"SMILES Column {smiles_col} not found in dataframe")

    for i, chunk in enumerate(chunks):

        smi_list = chunk[smiles_col].tolist()
        mol_list = [Chem.MolFromSmiles(smi) for smi in smi_list]
        legends = []
        if legend_columns:
            for col in legend_columns:
                legends.append(chunk[col].tolist())
            # zip each one into a string
            legends = [
                " ".join([f"{col}: {val}" for col, val in zip(legend_columns, leg)])
                for leg in zip(*legends)
            ]
        else:
            legends = ["" for _ in range(len(mol_list))]

        # make high quality image
        img = Draw.MolsToGridImage(
            mol_list,
            molsPerRow=mols_per_row,
            legends=legends,
            subImgSize=(150, 150),
            useSVG=True,
            returnPNG=True,
        )
        if rdkit_draw_in_ipython_mode():
            svg = img.data
        else:
            svg = img
        fname_svg = str(img_path / f"mols_{i}.svg")
        with open(fname_svg, "w") as outf:
            outf.write(svg)
        svgnames.append(fname_svg)

        # convert svg to png
        fname_png = str(img_path / f"mols_{i}.png")
        cairosvg.svg2png(url=fname_svg, write_to=fname_png)
        pngnames.append(fname_png)

    create_presentation(pngnames, output_file)
